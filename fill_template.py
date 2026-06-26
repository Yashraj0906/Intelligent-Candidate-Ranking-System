"""
Fill in the Redrob Idea Submission Template with our actual project content.
Preserves the original template design (backgrounds, colors, fonts).
"""
import sys, copy
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation(r'Idea Submission Template _ Redrob.pptx')

WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT = RGBColor(0xcc, 0xcc, 0xcc)
ACCENT = RGBColor(0x8b, 0x5c, 0xf6)
GREEN = RGBColor(0x22, 0xc5, 0x5e)


def clear_and_write(text_frame, lines, font_size=Pt(120000), color=WHITE):
    """Clear existing text and write new content, preserving first run's formatting."""
    # Get formatting reference from first run
    ref_font = None
    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
        ref_font = text_frame.paragraphs[0].runs[0].font

    # Clear all paragraphs
    for para in text_frame.paragraphs:
        para.clear()

    # Write first line into existing first paragraph
    if lines:
        first = lines[0]
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = first[0] if isinstance(first, tuple) else first
        sz = first[1] if isinstance(first, tuple) and len(first) > 1 else Pt(13)
        cl = first[2] if isinstance(first, tuple) and len(first) > 2 else WHITE
        bd = first[3] if isinstance(first, tuple) and len(first) > 3 else False
        run.font.size = sz
        run.font.color.rgb = cl
        run.font.bold = bd
        if ref_font and ref_font.name:
            run.font.name = ref_font.name

    # Add remaining lines
    for line in lines[1:]:
        p = text_frame.add_paragraph()
        run = p.add_run()
        if isinstance(line, tuple):
            run.text = line[0]
            run.font.size = line[1] if len(line) > 1 else Pt(13)
            run.font.color.rgb = line[2] if len(line) > 2 else WHITE
            run.font.bold = line[3] if len(line) > 3 else False
        else:
            run.text = line
            run.font.size = Pt(13)
            run.font.color.rgb = WHITE
        if ref_font and ref_font.name:
            run.font.name = ref_font.name
        p.space_before = Pt(2)


def add_content_box(slide, left, top, width, height, lines):
    """Add a new text box with content lines."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(2)

        run = p.add_run()
        if isinstance(line, tuple):
            run.text = line[0]
            run.font.size = line[1] if len(line) > 1 else Pt(12)
            run.font.color.rgb = line[2] if len(line) > 2 else WHITE
            run.font.bold = line[3] if len(line) > 3 else False
        else:
            run.text = line
            run.font.size = Pt(12)
            run.font.color.rgb = WHITE
        run.font.name = "Calibri"


slides = list(prs.slides)

# ============================================================
# SLIDE 1: Title slide - fill in team info
# ============================================================
s = slides[0]
tf = s.shapes[1].text_frame  # The text box with "Team Name :"
clear_and_write(tf, [
    ("Team Name:  InternDedo", Pt(16), WHITE, True),
    ("", Pt(6)),
    ("Problem Statement:  Rank top 100 from 100K candidates", Pt(14), WHITE, False),
    ("for a Senior AI Engineer role at Redrob AI", Pt(13), LIGHT),
    ("", Pt(6)),
    ("Team Leader:  Yashraj Kumar", Pt(14), WHITE, True),
    ("IIIT Nagpur  |  bt23ece018@iiitn.ac.in", Pt(12), LIGHT),
])

# ============================================================
# SLIDE 2: Solution Overview
# ============================================================
s = slides[1]
tf = s.shapes[1].text_frame  # Content text box
clear_and_write(tf, [
    ("What is our proposed solution?", Pt(14), ACCENT, True),
    ("", Pt(4)),
    ("I built a 5-stage AI pipeline that doesn't just match keywords --", Pt(12), WHITE),
    ("it actually understands who'd be a good Senior AI Engineer.", Pt(12), WHITE),
    ("", Pt(4)),
    ("The core idea: combine two types of understanding:", Pt(12), WHITE, True),
    ("  1. Semantic understanding via BGE embeddings", Pt(11), LIGHT),
    ("     (what the candidate's career is really about)", Pt(11), LIGHT),
    ("  2. Structured signal analysis", Pt(11), LIGHT),
    ("     (years of exp, skill depth, production work, title trajectory)", Pt(11), LIGHT),
    ("  3. Deep reranking via cross-encoder that reads JD + candidate together", Pt(11), LIGHT),
    ("", Pt(4)),
    ("Then I fuse everything using Reciprocal Rank Fusion --", Pt(12), WHITE),
    ("no manual weight tuning needed for the retrieval step.", Pt(12), WHITE),
    ("", Pt(6)),
    ("What differentiates our approach?", Pt(14), ACCENT, True),
    ("", Pt(4)),
    ("  -> Bi-encoder + Cross-encoder (like Google Search)", Pt(11), LIGHT),
    ("  -> 7 honeypot heuristics that catch fake profiles", Pt(11), LIGHT),
    ("  -> Anti-keyword-stuffer logic", Pt(11), LIGHT),
    ("  -> Consulting vs product company detection", Pt(11), LIGHT),
    ("  -> Skill authenticity scoring", Pt(11), LIGHT),
    ("  -> Every ranking has fact-based reasoning, zero hallucination", Pt(11), LIGHT),
])

# ============================================================
# SLIDE 3: JD Understanding
# ============================================================
s = slides[2]
tf = s.shapes[1].text_frame
clear_and_write(tf, [
    ("Key requirements I extracted from the JD:", Pt(14), ACCENT, True),
    ("", Pt(4)),
    ("Must-have skill buckets (I grouped related skills):", Pt(12), WHITE, True),
    ("  1. Embeddings/Retrieval: sentence-transformers, FAISS, Pinecone, ChromaDB", Pt(11), LIGHT),
    ("  2. Vector DB/Search: Elasticsearch, Milvus, Weaviate, Vespa", Pt(11), LIGHT),
    ("  3. Python/Coding: Python, PyTorch, TensorFlow, NumPy", Pt(11), LIGHT),
    ("  4. Ranking/Evaluation: NDCG, MRR, learning-to-rank, reranking", Pt(11), LIGHT),
    ("", Pt(4)),
    ("Nice-to-have: LLM fine-tuning, MLOps, open-source contributions", Pt(11), LIGHT),
    ("", Pt(6)),
    ("How I evaluate fit beyond keywords:", Pt(14), ACCENT, True),
    ("", Pt(4)),
    ("I don't just check if someone listed 'PyTorch' -- I check:", Pt(12), WHITE),
    ("  - Career trajectory: is the title path heading toward AI/ML?", Pt(11), LIGHT),
    ("  - Skill authenticity: are listed skills mentioned in job descriptions?", Pt(11), LIGHT),
    ("  - Production signals: words like 'deployed', 'serving', 'latency'", Pt(11), LIGHT),
    ("  - Experience sweet spot: Gaussian scoring centered at ~7 years", Pt(11), LIGHT),
    ("  - Consulting vs Product: entire career at TCS/Infosys = penalty", Pt(11), LIGHT),
    ("  - Behavioral: open to work, recruiter response rate, GitHub activity", Pt(11), LIGHT),
    ("  - Semantic similarity: BGE embedding distance between JD and resume", Pt(11), LIGHT),
])

# ============================================================
# SLIDE 4: Ranking Methodology
# ============================================================
s = slides[3]
tf = s.shapes[1].text_frame
clear_and_write(tf, [
    ("5-Stage Pipeline:", Pt(14), ACCENT, True),
    ("", Pt(4)),
    ("Stage 1 - Hard Filters:", Pt(12), GREEN, True),
    ("  Remove honeypots (528 found) + experience band (2-20 yrs)", Pt(11), LIGHT),
    ("  100K -> ~92K viable candidates", Pt(11), LIGHT),
    ("", Pt(4)),
    ("Stage 2 - Hybrid Retrieval (92K -> 2,000):", Pt(12), GREEN, True),
    ("  Dense: BGE-small-en-v1.5 cosine similarity (semantic)", Pt(11), LIGHT),
    ("  Sparse: Feature-based scoring (keyword/skill proxy)", Pt(11), LIGHT),
    ("  Fusion: RRF = 1/(k+rank_dense) + 1/(k+rank_sparse)", Pt(11), LIGHT),
    ("", Pt(4)),
    ("Stage 3 - Multi-Signal Reranker (2,000 -> 200):", Pt(12), GREEN, True),
    ("  Career Fit 35% | Skills Fit 25% | Behavioral 20% | Logistics 10% | Background 10%", Pt(11), LIGHT),
    ("", Pt(4)),
    ("Stage 4 - Cross-Encoder Reranker (200 -> 100):", Pt(12), GREEN, True),
    ("  ms-marco-MiniLM-L-6-v2 reads (JD, candidate) together", Pt(11), LIGHT),
    ("  Final = 70% multi-signal + 30% cross-encoder", Pt(11), LIGHT),
    ("", Pt(4)),
    ("Stage 5 - Reasoning Generation:", Pt(12), GREEN, True),
    ("  Fact-based, per-candidate reasoning from actual data -- no LLM", Pt(11), LIGHT),
])

# ============================================================
# SLIDE 5: Explainability & Data Validation
# ============================================================
s = slides[4]
tf = s.shapes[1].text_frame
clear_and_write(tf, [
    ("How are ranking decisions explained?", Pt(14), ACCENT, True),
    ("", Pt(3)),
    ("Every candidate gets a fact-based reasoning string. I wrote a", Pt(11), WHITE),
    ("rule-based generator that references actual skills, companies, titles.", Pt(11), WHITE),
    ("Example: 'Senior ML Engineer at Razorpay with 7.2 yrs experience;", Pt(11), LIGHT),
    ("covers 3/4 must-have skill areas: embeddings, vector DB, Python'", Pt(11), LIGHT),
    ("", Pt(3)),
    ("How I prevent hallucinations:", Pt(12), ACCENT, True),
    ("No LLM is used for reasoning. Everything is template-based from", Pt(11), WHITE),
    ("extracted features. If a skill isn't in the profile, it's never mentioned.", Pt(11), WHITE),
    ("", Pt(3)),
    ("Handling suspicious profiles (7 honeypot heuristics):", Pt(12), ACCENT, True),
    ("  H1: Experience vs career history mismatch", Pt(10), LIGHT),
    ("  H2: Expert skills with 0 months usage", Pt(10), LIGHT),
    ("  H3: Too many expert skills + no assessments taken", Pt(10), LIGHT),
    ("  H4: All skills at exact same proficiency (synthetic pattern)", Pt(10), LIGHT),
    ("  H5: Title-description mismatch", Pt(10), LIGHT),
    ("  H6: Impossible tenure (20+ yrs at one company)", Pt(10), LIGHT),
    ("  H7: Endorsement anomaly (high endorsements, low duration)", Pt(10), LIGHT),
    ("", Pt(3)),
    ("Result: 528 honeypots caught, 0 in final top 100", Pt(12), GREEN, True),
])

# ============================================================
# SLIDE 6: End-to-End Workflow
# ============================================================
s = slides[5]
tf = s.shapes[0].text_frame  # This slide only has one text box
clear_and_write(tf, [
    ("End-to-End Workflow", Pt(16), WHITE, True),
    ("", Pt(4)),
    ("OFFLINE (run once, ~20 min):", Pt(13), RGBColor(0xfb, 0x92, 0x3c), True),
    ("  python precompute.py --candidates candidates.jsonl", Pt(11), LIGHT),
    ("  1. Load 100K candidates from JSONL", Pt(11), LIGHT),
    ("  2. Extract 50+ structured features per candidate", Pt(11), LIGHT),
    ("  3. Run 7 honeypot heuristics -> flag 528 suspicious", Pt(11), LIGHT),
    ("  4. Pre-filter via feature scoring -> top 5,000 shortlist", Pt(11), LIGHT),
    ("  5. Compute BGE-small embeddings for 5,000 (5000 x 384)", Pt(11), LIGHT),
    ("  6. Save all artifacts to data/ directory", Pt(11), LIGHT),
    ("", Pt(4)),
    ("ONLINE (87 seconds, CPU only, no network):", Pt(13), GREEN, True),
    ("  python rank.py --candidates candidates.jsonl --out submission.csv", Pt(11), LIGHT),
    ("  1. Load pre-computed artifacts (~3s)", Pt(11), LIGHT),
    ("  2. Hybrid retrieval: Dense + Sparse + RRF -> top 2,000", Pt(11), LIGHT),
    ("  3. Multi-signal reranking (5 components) -> top 200", Pt(11), LIGHT),
    ("  4. Cross-encoder reranking (ms-marco-MiniLM) -> top 100 (~84s)", Pt(11), LIGHT),
    ("  5. Generate reasoning + write submission.csv", Pt(11), LIGHT),
])

# ============================================================
# SLIDE 7: System Architecture
# ============================================================
s = slides[6]
# Keep the title, add architecture below
add_content_box(s, 572589, 1300000, 8000000, 3600000, [
    ("100K candidates  +  JD (Senior AI Engineer)", Pt(13), ACCENT, True),
    ("       |                              |", Pt(10), LIGHT),
    ("[Feature Extraction]          [JD Embedding: BGE-small]", Pt(12), WHITE, True),
    ("       |", Pt(10), LIGHT),
    ("[Honeypot Detection] -- 528 flagged & removed", Pt(12), RGBColor(0xfb, 0x92, 0x3c), True),
    ("       |", Pt(10), LIGHT),
    ("[Pre-Filter] -> top 5,000 by feature score", Pt(12), WHITE),
    ("       |", Pt(10), LIGHT),
    ("[BGE Embeddings] -> 5000 x 384 dense vectors", Pt(12), WHITE),
    ("       |                              |", Pt(10), LIGHT),
    ("[Dense Retrieval]          [Sparse Retrieval]", Pt(12), WHITE),
    ("       \\                        /", Pt(10), LIGHT),
    ("  [Reciprocal Rank Fusion] -> top 2,000", Pt(12), RGBColor(0x3b, 0x82, 0xf6), True),
    ("             |", Pt(10), LIGHT),
    ("  [Multi-Signal Reranker] -> top 200", Pt(12), RGBColor(0x3b, 0x82, 0xf6), True),
    ("  Career(35%) + Skills(25%) + Behavioral(20%) + Logistics(10%) + BG(10%)", Pt(10), LIGHT),
    ("             |", Pt(10), LIGHT),
    ("  [Cross-Encoder Reranker] -> top 100", Pt(12), GREEN, True),
    ("  ms-marco-MiniLM-L-6-v2 | 70% multi-signal + 30% cross-encoder", Pt(10), LIGHT),
    ("             |", Pt(10), LIGHT),
    ("  [Reasoning Generator] -> submission.csv", Pt(12), ACCENT, True),
])

# ============================================================
# SLIDE 8: Results & Performance
# ============================================================
s = slides[7]
tf = s.shapes[1].text_frame
clear_and_write(tf, [
    ("Key results:", Pt(14), ACCENT, True),
    ("", Pt(3)),
    ("  Total candidates processed:    100,000", Pt(12), WHITE),
    ("  Honeypots detected & filtered: 528", Pt(12), WHITE),
    ("  Honeypots in top 100:          0  (limit: <=10)", Pt(12), GREEN, True),
    ("  Keyword stuffers in top 100:   0", Pt(12), GREEN, True),
    ("  AI/ML titles in top 100:       90/100", Pt(12), WHITE),
    ("  Validation:                    PASSED", Pt(12), GREEN, True),
    ("", Pt(3)),
    ("Runtime & compute:", Pt(14), ACCENT, True),
    ("  Ranking runtime:    87 seconds (limit: 300s)", Pt(12), WHITE),
    ("  Headroom:           3.5 minutes spare", Pt(12), GREEN),
    ("  CPU only:           Yes (no GPU)", Pt(12), WHITE),
    ("  Network:            None (fully offline)", Pt(12), WHITE),
    ("", Pt(3)),
    ("Live demo tested (HuggingFace Spaces):", Pt(12), ACCENT, True),
    ("  AI Engineer scored 0.8353 -> ranked #1", Pt(11), GREEN),
    ("  Ops Manager scored 0.2614 -> correctly ranked lower", Pt(11), LIGHT),
    ("  Honeypot profile -> correctly filtered out", Pt(11), RGBColor(0xfb, 0x92, 0x3c)),
])

# ============================================================
# SLIDE 9: Technologies Used
# ============================================================
s = slides[8]
tf = s.shapes[1].text_frame
clear_and_write(tf, [
    ("Why I chose each technology:", Pt(14), ACCENT, True),
    ("", Pt(3)),
    ("Bi-Encoder: BAAI/bge-small-en-v1.5 (33M params)", Pt(12), WHITE, True),
    ("  Top-ranked on MTEB for its size. 384-dim. 3x faster than bge-base on CPU.", Pt(11), LIGHT),
    ("", Pt(3)),
    ("Cross-Encoder: ms-marco-MiniLM-L-6-v2 (22M params)", Pt(12), WHITE, True),
    ("  Trained on MS-MARCO for passage ranking. Reads JD + candidate", Pt(11), LIGHT),
    ("  together via full self-attention. Catches subtle matches.", Pt(11), LIGHT),
    ("", Pt(3)),
    ("Reciprocal Rank Fusion (RRF)", Pt(12), WHITE, True),
    ("  Production-grade rank merging used by Elasticsearch 8.x and Vespa.", Pt(11), LIGHT),
    ("  No weight tuning needed.", Pt(11), LIGHT),
    ("", Pt(3)),
    ("Python + NumPy + sentence-transformers", Pt(12), WHITE, True),
    ("  Pure Python. NumPy for vectorized ops on 100K candidates.", Pt(11), LIGHT),
    ("  Pickle/NPY for fast artifact serialization.", Pt(11), LIGHT),
    ("", Pt(3)),
    ("Gradio on HuggingFace Spaces", Pt(12), WHITE, True),
    ("  Free tier, Python-native, instant deployment for live demo.", Pt(11), LIGHT),
])

# ============================================================
# SLIDE 10: Submission Assets
# ============================================================
s = slides[9]
tf = s.shapes[0].text_frame
clear_and_write(tf, [
    ("Submission Assets", Pt(16), WHITE, True),
    ("", Pt(6)),
    ("GitHub Repository:", Pt(13), ACCENT, True),
    ("  github.com/Yashraj0906/Intelligent-Candidate-Ranking-System", Pt(11), RGBColor(0x3b, 0x82, 0xf6)),
    ("  Full source code, README with Mermaid architecture diagrams", Pt(11), LIGHT),
    ("", Pt(4)),
    ("Live Demo (HuggingFace Spaces):", Pt(13), ACCENT, True),
    ("  huggingface.co/spaces/yashrajkumar623/Intelligent-Candidate-Ranking-System", Pt(11), RGBColor(0x3b, 0x82, 0xf6)),
    ("  Paste candidate JSON -> get instant rankings with reasoning", Pt(11), LIGHT),
    ("", Pt(4)),
    ("Reproduce the ranking:", Pt(13), ACCENT, True),
    ("  python rank.py --candidates candidates.jsonl --out submission.csv", Pt(11), GREEN),
    ("", Pt(4)),
    ("AI Tools: Claude (via Antigravity IDE) for architecture discussion", Pt(12), LIGHT),
    ("and code generation. No candidate data sent to any hosted LLM.", Pt(12), LIGHT),
    ("Ranking runs 100% offline on CPU.", Pt(12), WHITE, True),
])

# Save
output_path = r'InternDedo_Submission_Final.pptx'
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
