"""
Generate submission PPT for Redrob Hackathon.
Fills in the Idea Submission Template with our actual project details.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color scheme
BG_DARK = RGBColor(0x0f, 0x17, 0x2a)
ACCENT_BLUE = RGBColor(0x3b, 0x82, 0xf6)
ACCENT_PURPLE = RGBColor(0x8b, 0x5c, 0xf6)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xcc, 0xcc, 0xcc)
GREEN = RGBColor(0x22, 0xc5, 0x5e)
ORANGE = RGBColor(0xfb, 0x92, 0x3c)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=16, bold=False, color=WHITE, space_before=6, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.space_before = Pt(space_before)
    return p

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)

add_text_box(slide, 1.5, 1.0, 10, 1.2,
             "Intelligent Candidate Discovery & Ranking System",
             font_size=36, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 2.5, 10, 0.6,
             "Redrob Hackathon: India Runs Data & AI Challenge 2026",
             font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

tf = add_text_box(slide, 3, 3.8, 7, 2.5, "", font_size=18, alignment=PP_ALIGN.CENTER)
add_para(tf, "Team Name: InternDedo", font_size=22, bold=True, color=WHITE, space_before=12)
add_para(tf, "", font_size=10, space_before=8)
add_para(tf, "Problem Statement: Rank top 100 candidates from 100K profiles", font_size=18, color=LIGHT_GRAY, space_before=8)
add_para(tf, "for a Senior AI Engineer role at Redrob AI", font_size=18, color=LIGHT_GRAY, space_before=4)
add_para(tf, "", font_size=10, space_before=12)
add_para(tf, "Team Leader: Yashraj Kumar", font_size=20, bold=True, color=WHITE, space_before=12)
add_para(tf, "bt23ece018@iiitn.ac.in | IIIT Nagpur", font_size=16, color=LIGHT_GRAY, space_before=4)

# ============================================================
# SLIDE 2: Solution Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Solution Overview", font_size=32, bold=True, color=ACCENT_BLUE)

tf = add_text_box(slide, 0.8, 1.4, 5.5, 5, "", font_size=14)
add_para(tf, "What is our proposed solution?", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=0)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "I built a 5-stage AI pipeline that doesn't just match keywords — it actually understands who would be a good Senior AI Engineer.", font_size=16, color=WHITE, space_before=6)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "The core idea: combine two types of understanding —", font_size=16, color=WHITE, space_before=4)
add_para(tf, "  1. Semantic understanding via BGE embeddings (what the candidate's career is really about)", font_size=15, color=LIGHT_GRAY, space_before=6)
add_para(tf, "  2. Structured signal analysis (years of experience, skill depth, production work, title trajectory)", font_size=15, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  3. Deep reranking via cross-encoder that reads JD + candidate together", font_size=15, color=LIGHT_GRAY, space_before=4)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "Then I fuse everything using Reciprocal Rank Fusion — no manual weight tuning needed for the retrieval step.", font_size=16, color=WHITE, space_before=6)

tf = add_text_box(slide, 7, 1.4, 5.5, 5, "", font_size=14)
add_para(tf, "What differentiates our approach?", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=0)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "Most systems do keyword matching or basic TF-IDF. Mine is different:", font_size=16, color=WHITE, space_before=6)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "  -> Bi-encoder + Cross-encoder architecture (like how Google Search works)", font_size=15, color=LIGHT_GRAY, space_before=6)
add_para(tf, "  -> 7 honeypot heuristics that catch fake/impossible profiles", font_size=15, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  -> Anti-keyword-stuffer logic (catches Marketing Managers with inflated AI skill lists)", font_size=15, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  -> Consulting vs product company detection", font_size=15, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  -> Skill authenticity scoring (do skills actually appear in career descriptions?)", font_size=15, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  -> Every ranking comes with fact-based reasoning, zero hallucination", font_size=15, color=LIGHT_GRAY, space_before=4)

# ============================================================
# SLIDE 3: JD Understanding
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "JD Understanding & Candidate Evaluation", font_size=32, bold=True, color=ACCENT_BLUE)

tf = add_text_box(slide, 0.8, 1.4, 5.5, 5.5, "", font_size=14)
add_para(tf, "Key requirements I extracted from the JD:", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=0)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "Must-have skill buckets (I grouped related skills):", font_size=16, bold=True, color=WHITE, space_before=6)
add_para(tf, "  1. Embeddings/Retrieval: sentence-transformers, FAISS, Pinecone, ChromaDB", font_size=14, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  2. Vector DB/Search: Elasticsearch, Milvus, Weaviate, Vespa", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  3. Python/Coding: Python, PyTorch, TensorFlow, NumPy", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  4. Ranking/Evaluation: NDCG, MRR, learning-to-rank, reranking", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "Nice-to-have:", font_size=16, bold=True, color=WHITE, space_before=6)
add_para(tf, "  - LLM fine-tuning (LoRA, RLHF, DPO)", font_size=14, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  - MLOps (Kubernetes, MLflow, Docker)", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  - Open-source contributions", font_size=14, color=LIGHT_GRAY, space_before=3)

tf = add_text_box(slide, 7, 1.4, 5.5, 5.5, "", font_size=14)
add_para(tf, "How I evaluate candidate fit beyond keywords:", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=0)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "I don't just check if someone listed 'PyTorch' — I check:", font_size=16, color=WHITE, space_before=6)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "  Career trajectory — is the title path heading toward AI/ML?", font_size=15, color=LIGHT_GRAY, space_before=6)
add_para(tf, "  Skill authenticity — are listed skills actually mentioned in job descriptions?", font_size=15, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  Production signals — words like 'deployed', 'serving', 'latency' in descriptions", font_size=15, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  Experience sweet spot — Gaussian scoring centered at ~7 years", font_size=15, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  Consulting vs Product — entire career at TCS/Infosys = penalty", font_size=15, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  Behavioral signals — open to work, recruiter response rate, GitHub activity", font_size=15, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  Semantic similarity — BGE embedding distance between JD and resume text", font_size=15, color=LIGHT_GRAY, space_before=4)

# ============================================================
# SLIDE 4: Ranking Methodology
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Ranking Methodology", font_size=32, bold=True, color=ACCENT_BLUE)

tf = add_text_box(slide, 0.8, 1.3, 11.5, 6, "", font_size=14)
add_para(tf, "How the system retrieves, scores, and ranks:", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=0)
add_para(tf, "", font_size=4, space_before=4)

add_para(tf, "Stage 1 - Hard Filters: Remove honeypots (528 found) + experience band filter (2-20 yrs)", font_size=15, bold=True, color=GREEN, space_before=8)
add_para(tf, "  Result: 100K -> ~92K viable candidates", font_size=14, color=LIGHT_GRAY, space_before=3)

add_para(tf, "Stage 2 - Hybrid Retrieval (92K -> 2,000):", font_size=15, bold=True, color=GREEN, space_before=10)
add_para(tf, "  Dense: BGE-small-en-v1.5 cosine similarity (semantic understanding)", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  Sparse: Feature-based scoring (keyword/skill matching proxy)", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  Fusion: RRF = 1/(k+rank_dense) + 1/(k+rank_sparse) — no weight tuning needed", font_size=14, color=LIGHT_GRAY, space_before=3)

add_para(tf, "Stage 3 - Multi-Signal Reranker (2,000 -> 200):", font_size=15, bold=True, color=GREEN, space_before=10)
add_para(tf, "  Career Fit (35%) | Skills Fit (25%) | Behavioral (20%) | Logistics (10%) | Background (10%)", font_size=14, color=LIGHT_GRAY, space_before=3)

add_para(tf, "Stage 4 - Cross-Encoder Reranker (200 -> 100):", font_size=15, bold=True, color=GREEN, space_before=10)
add_para(tf, "  ms-marco-MiniLM-L-6-v2 reads (JD, candidate) pairs together via full self-attention", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  Final score = 70% multi-signal composite + 30% cross-encoder normalized", font_size=14, color=LIGHT_GRAY, space_before=3)

add_para(tf, "Stage 5 - Reasoning Generation:", font_size=15, bold=True, color=GREEN, space_before=10)
add_para(tf, "  Fact-based, per-candidate reasoning from actual profile data — no LLM, no hallucination", font_size=14, color=LIGHT_GRAY, space_before=3)

# ============================================================
# SLIDE 5: Explainability & Data Validation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Explainability & Data Validation", font_size=32, bold=True, color=ACCENT_BLUE)

tf = add_text_box(slide, 0.8, 1.4, 5.5, 5.5, "", font_size=14)
add_para(tf, "How are ranking decisions explained?", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=0)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "Every candidate gets a specific, fact-based reasoning string. I wrote a rule-based reasoning generator that:", font_size=15, color=WHITE, space_before=6)
add_para(tf, "", font_size=4, space_before=2)
add_para(tf, "  - References actual skills found in the profile", font_size=14, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  - Mentions real companies and titles", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  - Cites specific must-have skill coverage (e.g., '3/4 buckets')", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  - Flags concerns like 'all consulting' or 'above preferred experience range'", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "How I prevent hallucinations:", font_size=18, bold=True, color=ORANGE, space_before=8)
add_para(tf, "  No LLM is used for reasoning. Everything is template-based from extracted features. If a skill isn't in the profile, it's never mentioned.", font_size=14, color=LIGHT_GRAY, space_before=4)

tf = add_text_box(slide, 7, 1.4, 5.5, 5.5, "", font_size=14)
add_para(tf, "Handling suspicious profiles:", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=0)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "7 independent honeypot heuristics:", font_size=16, bold=True, color=WHITE, space_before=6)
add_para(tf, "  H1: Experience vs career history mismatch", font_size=14, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  H2: Expert skills with 0 months usage", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  H3: Too many expert skills + no assessments", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  H4: All skills at exact same proficiency (synthetic)", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  H5: Title-description mismatch", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  H6: Impossible tenure (20+ yrs at one company)", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  H7: Endorsement anomaly", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "Result: 528 honeypots caught, 0 in final top 100", font_size=16, bold=True, color=GREEN, space_before=6)
add_para(tf, "", font_size=4, space_before=2)
add_para(tf, "Anti-stuffer: Non-tech titles + high AI keyword density + low skill mention rate in career descriptions = score reduced to 15%", font_size=14, color=LIGHT_GRAY, space_before=4)

# ============================================================
# SLIDE 6: End-to-End Workflow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "End-to-End Workflow", font_size=32, bold=True, color=ACCENT_BLUE)

tf = add_text_box(slide, 0.8, 1.3, 11.5, 5.5, "", font_size=14)
add_para(tf, "Complete workflow from JD input to ranked output:", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=0)
add_para(tf, "", font_size=8, space_before=4)

add_para(tf, "OFFLINE (run once, ~20 min):", font_size=18, bold=True, color=ORANGE, space_before=8)
add_para(tf, "  python precompute.py --candidates candidates.jsonl", font_size=14, color=LIGHT_GRAY, space_before=6)
add_para(tf, "    1. Load 100K candidates from JSONL", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "    2. Extract 50+ structured features per candidate (titles, skills, production signals...)", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "    3. Run 7 honeypot heuristics -> flag 528 suspicious profiles", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "    4. Pre-filter via feature scoring -> top 5,000 shortlist", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "    5. Compute BGE-small-en-v1.5 embeddings for 5,000 shortlisted (5000 x 384)", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "    6. Compute JD embedding (1 x 384)", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "    7. Save all artifacts to data/ directory (pickle + numpy)", font_size=14, color=LIGHT_GRAY, space_before=3)

add_para(tf, "", font_size=6, space_before=6)
add_para(tf, "ONLINE (87 seconds, CPU only, no network):", font_size=18, bold=True, color=GREEN, space_before=8)
add_para(tf, "  python rank.py --candidates candidates.jsonl --out submission.csv", font_size=14, color=LIGHT_GRAY, space_before=6)
add_para(tf, "    1. Load pre-computed artifacts (~3s)", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "    2. Hybrid retrieval: Dense + Sparse + RRF -> top 2,000", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "    3. Multi-signal reranking (5 components) -> top 200", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "    4. Cross-encoder reranking (ms-marco-MiniLM) -> top 100 (~84s)", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "    5. Generate reasoning + write submission.csv", font_size=14, color=LIGHT_GRAY, space_before=3)

# ============================================================
# SLIDE 7: System Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "System Architecture", font_size=32, bold=True, color=ACCENT_BLUE)

# Architecture as a text-based diagram since we can't embed mermaid in pptx
tf = add_text_box(slide, 0.5, 1.2, 12.3, 5.8, "", font_size=13)
add_para(tf, "candidates.jsonl (100K)     +     Job Description (Senior AI Engineer)", font_size=16, bold=True, color=ACCENT_PURPLE, space_before=0)
add_para(tf, "         |                                    |", font_size=13, color=LIGHT_GRAY, space_before=2)
add_para(tf, "         v                                    v", font_size=13, color=LIGHT_GRAY, space_before=2)
add_para(tf, "  [Feature Extraction]                [JD Embedding]", font_size=15, color=WHITE, space_before=2)
add_para(tf, "  100K -> 50+ features each           BGE-small (1x384)", font_size=13, color=LIGHT_GRAY, space_before=2)
add_para(tf, "         |", font_size=13, color=LIGHT_GRAY, space_before=2)
add_para(tf, "  [Honeypot Detection] -- 528 flagged, removed", font_size=15, color=ORANGE, space_before=2)
add_para(tf, "         |", font_size=13, color=LIGHT_GRAY, space_before=2)
add_para(tf, "  [Pre-Filter] -> top 5,000 by feature score", font_size=15, color=WHITE, space_before=2)
add_para(tf, "         |", font_size=13, color=LIGHT_GRAY, space_before=2)
add_para(tf, "  [BGE Embeddings] -> 5000 x 384 dense vectors", font_size=15, color=WHITE, space_before=2)
add_para(tf, "         |                                    |", font_size=13, color=LIGHT_GRAY, space_before=2)
add_para(tf, "  [Dense Retrieval]              [Sparse Retrieval]", font_size=15, color=WHITE, space_before=2)
add_para(tf, "         \\                              /", font_size=13, color=LIGHT_GRAY, space_before=2)
add_para(tf, "          [Reciprocal Rank Fusion] -> top 2,000", font_size=15, bold=True, color=ACCENT_BLUE, space_before=2)
add_para(tf, "                     |", font_size=13, color=LIGHT_GRAY, space_before=2)
add_para(tf, "          [Multi-Signal Reranker] -> top 200", font_size=15, bold=True, color=ACCENT_BLUE, space_before=2)
add_para(tf, "          Career(35%) + Skills(25%) + Behavioral(20%) + Logistics(10%) + Background(10%)", font_size=12, color=LIGHT_GRAY, space_before=2)
add_para(tf, "                     |", font_size=13, color=LIGHT_GRAY, space_before=2)
add_para(tf, "          [Cross-Encoder Reranker] -> top 100", font_size=15, bold=True, color=GREEN, space_before=2)
add_para(tf, "          ms-marco-MiniLM-L-6-v2  |  70% multi-signal + 30% cross-encoder", font_size=12, color=LIGHT_GRAY, space_before=2)
add_para(tf, "                     |", font_size=13, color=LIGHT_GRAY, space_before=2)
add_para(tf, "          [Reasoning Generator] -> submission.csv", font_size=15, bold=True, color=ACCENT_PURPLE, space_before=2)

# ============================================================
# SLIDE 8: Results & Performance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Results & Performance", font_size=32, bold=True, color=ACCENT_BLUE)

tf = add_text_box(slide, 0.8, 1.3, 5.5, 5.5, "", font_size=14)
add_para(tf, "Ranking quality results:", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=0)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "  Total candidates processed:     100,000", font_size=16, color=WHITE, space_before=6)
add_para(tf, "  Honeypots detected & filtered:  528", font_size=16, color=WHITE, space_before=4)
add_para(tf, "  Honeypots in top 100:           0  (limit: <=10)", font_size=16, bold=True, color=GREEN, space_before=4)
add_para(tf, "  Keyword stuffers in top 100:    0", font_size=16, bold=True, color=GREEN, space_before=4)
add_para(tf, "  AI/ML titles in top 100:        90/100", font_size=16, color=WHITE, space_before=4)
add_para(tf, "  Non-tech titles in top 100:     0", font_size=16, bold=True, color=GREEN, space_before=4)
add_para(tf, "  Validation:                     PASSED", font_size=16, bold=True, color=GREEN, space_before=4)
add_para(tf, "", font_size=8, space_before=4)
add_para(tf, "Top 5 ranked candidates:", font_size=18, bold=True, color=WHITE, space_before=8)
add_para(tf, "  1. Sr ML Engineer @ Zomato (7.2yr, CE:0.99)", font_size=14, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  2. Staff ML Engineer @ Yellow.ai (8.6yr, CE:0.94)", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  3. Sr AI Engineer @ Apple (5.9yr, CE:1.00)", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  4. Staff ML Engineer @ Paytm (7.0yr, CE:0.92)", font_size=14, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  5. Sr AI Engineer @ Netflix (7.8yr, CE:0.94)", font_size=14, color=LIGHT_GRAY, space_before=3)

tf = add_text_box(slide, 7, 1.3, 5.5, 5.5, "", font_size=14)
add_para(tf, "Runtime & compute constraints:", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=0)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "  Ranking runtime:      87 seconds", font_size=16, color=WHITE, space_before=6)
add_para(tf, "  Time limit:           5 minutes (300s)", font_size=16, color=LIGHT_GRAY, space_before=4)
add_para(tf, "  Headroom:             3.5 minutes spare", font_size=16, bold=True, color=GREEN, space_before=4)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "  CPU only:             Yes (no GPU needed)", font_size=16, color=WHITE, space_before=6)
add_para(tf, "  Network during rank:  None (fully offline)", font_size=16, color=WHITE, space_before=4)
add_para(tf, "  RAM used:             ~4 GB peak", font_size=16, color=WHITE, space_before=4)
add_para(tf, "", font_size=6, space_before=4)
add_para(tf, "Live demo tested on HuggingFace Spaces:", font_size=18, bold=True, color=WHITE, space_before=8)
add_para(tf, "  AI Engineer (score: 0.8353) ranked #1", font_size=15, color=GREEN, space_before=4)
add_para(tf, "  Operations Manager (score: 0.2614) ranked #2", font_size=15, color=LIGHT_GRAY, space_before=3)
add_para(tf, "  Honeypot profile correctly filtered out", font_size=15, color=ORANGE, space_before=3)

# ============================================================
# SLIDE 9: Technologies Used
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Technologies Used", font_size=32, bold=True, color=ACCENT_BLUE)

tf = add_text_box(slide, 0.8, 1.3, 11.5, 5.5, "", font_size=14)
add_para(tf, "Why I chose each technology:", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=0)
add_para(tf, "", font_size=6, space_before=4)

add_para(tf, "Bi-Encoder: BAAI/bge-small-en-v1.5 (33M params)", font_size=16, bold=True, color=WHITE, space_before=8)
add_para(tf, "  Top-ranked on MTEB benchmark for its size. 384-dim vectors. 3x faster than bge-base on CPU. I needed something that could encode 5000 candidates in minutes, not hours.", font_size=14, color=LIGHT_GRAY, space_before=3)

add_para(tf, "Cross-Encoder: ms-marco-MiniLM-L-6-v2 (22M params)", font_size=16, bold=True, color=WHITE, space_before=10)
add_para(tf, "  Trained on MS-MARCO passage ranking data. Unlike bi-encoders which encode separately, this reads JD + candidate together. More expensive but catches subtle matches bi-encoders miss.", font_size=14, color=LIGHT_GRAY, space_before=3)

add_para(tf, "Reciprocal Rank Fusion (RRF)", font_size=16, bold=True, color=WHITE, space_before=10)
add_para(tf, "  Used in production by Elasticsearch 8.x and Vespa. Merges dense + sparse rankings without needing to tune any weights. Formula: score = 1/(k+rank_dense) + 1/(k+rank_sparse)", font_size=14, color=LIGHT_GRAY, space_before=3)

add_para(tf, "Python + NumPy + sentence-transformers", font_size=16, bold=True, color=WHITE, space_before=10)
add_para(tf, "  Pure Python, no heavy infra needed. NumPy for vectorized ops on 100K candidates. Pickle/NPY for fast artifact serialization. Everything runs on a single machine.", font_size=14, color=LIGHT_GRAY, space_before=3)

add_para(tf, "Gradio on HuggingFace Spaces", font_size=16, bold=True, color=WHITE, space_before=10)
add_para(tf, "  Free tier, Python-native, instant deployment. Perfect for a demo sandbox — organizers can paste JSON and see results immediately.", font_size=14, color=LIGHT_GRAY, space_before=3)

# ============================================================
# SLIDE 10: Submission Assets
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Submission Assets", font_size=32, bold=True, color=ACCENT_BLUE)

tf = add_text_box(slide, 2, 1.5, 9, 5, "", font_size=14)
add_para(tf, "GitHub Repository:", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=0)
add_para(tf, "  https://github.com/Yashraj0906/Intelligent-Candidate-Ranking-System", font_size=16, color=ACCENT_BLUE, space_before=6)
add_para(tf, "  Full source code, README with Mermaid architecture diagrams, demo results", font_size=14, color=LIGHT_GRAY, space_before=3)

add_para(tf, "", font_size=10, space_before=8)
add_para(tf, "Live Demo (HuggingFace Spaces):", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=8)
add_para(tf, "  https://huggingface.co/spaces/yashrajkumar623/Intelligent-Candidate-Ranking-System", font_size=16, color=ACCENT_BLUE, space_before=6)
add_para(tf, "  Paste candidate JSON -> get instant rankings with per-candidate reasoning", font_size=14, color=LIGHT_GRAY, space_before=3)

add_para(tf, "", font_size=10, space_before=8)
add_para(tf, "Reproduce the ranking:", font_size=20, bold=True, color=ACCENT_PURPLE, space_before=8)
add_para(tf, "  python rank.py --candidates candidates.jsonl --out submission.csv", font_size=16, color=GREEN, space_before=6)
add_para(tf, "  Requires pre-computed artifacts in data/ directory (from precompute.py)", font_size=14, color=LIGHT_GRAY, space_before=3)

add_para(tf, "", font_size=10, space_before=8)
add_para(tf, "AI Tools Used: Claude (via Antigravity IDE) for architecture discussion and code generation.", font_size=15, color=LIGHT_GRAY, space_before=8)
add_para(tf, "No candidate data was sent to any hosted LLM. Ranking runs 100% offline on CPU.", font_size=15, bold=True, color=WHITE, space_before=3)

# ============================================================
# SLIDE 11: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.5, 2.0, 10, 1.2,
             "Thank You!", font_size=44, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

tf = add_text_box(slide, 3, 3.5, 7, 3, "", font_size=18, alignment=PP_ALIGN.CENTER)
add_para(tf, "Team InternDedo", font_size=28, bold=True, color=WHITE, space_before=0)
add_para(tf, "Yashraj Kumar | IIIT Nagpur", font_size=20, color=LIGHT_GRAY, space_before=12)
add_para(tf, "bt23ece018@iiitn.ac.in", font_size=18, color=LIGHT_GRAY, space_before=6)
add_para(tf, "", font_size=10, space_before=12)
add_para(tf, "Built with genuine curiosity about how search & ranking systems work", font_size=16, color=ACCENT_PURPLE, space_before=8)

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "InternDedo_Submission.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
